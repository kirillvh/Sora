"""markdown -> .pptx helper. You are NOT graded on slide formatting.

Supported markdown subset:

    # Deck Title          first level-1 heading -> title slide
    plain line            first plain line before any '##' -> subtitle
    ## Slide Title        starts a new content slide
    - bullet              bullet on the current slide
      - sub bullet        2+ spaces of indent -> second-level bullet
    plain line            under a '##': body paragraph on that slide

CLI:
    python -m tools.deck demo.md demo.pptx
"""
import argparse
import pathlib
import sys

from pptx import Presentation


def _parse(markdown: str):
    title, subtitle, slides = None, None, []
    current = None
    for raw in markdown.splitlines():
        line = raw.rstrip()
        if not line.strip():
            continue
        stripped = line.lstrip()
        if line.startswith("# ") and title is None and current is None:
            title = line[2:].strip()
        elif line.startswith("## "):
            current = {"title": line[3:].strip(), "items": []}
            slides.append(current)
        elif stripped.startswith("- "):
            level = 1 if (len(line) - len(stripped)) >= 2 else 0
            if current is None:
                current = {"title": "", "items": []}
                slides.append(current)
            current["items"].append((level, stripped[2:].strip()))
        else:
            if current is None:
                if subtitle is None:
                    subtitle = stripped
            else:
                current["items"].append((0, stripped))
    return title, subtitle, slides


def make_deck_from_markdown(markdown: str, out_path: str = "sora_deck.pptx") -> str:
    """Build a .pptx from the markdown subset above. Returns the saved path."""
    title, subtitle, slides = _parse(markdown)

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # title layout
    slide.shapes.title.text = title or "Untitled"
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle

    for spec in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
        slide.shapes.title.text = spec["title"]
        body = slide.placeholders[1].text_frame
        body.clear()
        first = True
        for level, text in spec["items"]:
            p = body.paragraphs[0] if first else body.add_paragraph()
            first = False
            p.text = text
            p.level = level

    pathlib.Path(out_path).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main():
    for _stream in (sys.stdout, sys.stderr):  # Windows consoles default to cp932/cp950
        if hasattr(_stream, "reconfigure"):
            _stream.reconfigure(encoding="utf-8", errors="replace")
    parser = argparse.ArgumentParser(description="markdown -> pptx")
    parser.add_argument("input_md")
    parser.add_argument("output_pptx")
    args = parser.parse_args()
    with open(args.input_md, encoding="utf-8") as f:
        path = make_deck_from_markdown(f.read(), args.output_pptx)
    print(f"Saved {path}")


if __name__ == "__main__":
    main()
